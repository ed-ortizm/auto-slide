"""Test pptx module"""
# import glob
import time
from configparser import ConfigParser, ExtendedInterpolation

from pptx import Presentation
from pptx.dml.color import RGBColor
# from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

###############################################################################
start_time = time.time()
###############################################################################
parser = ConfigParser(interpolation=ExtendedInterpolation())
CONFIG_FILE_NAME = "generate_slides.ini"
parser.read(f"{CONFIG_FILE_NAME}")
###############################################################################
presentation_name = parser.get("file", "presentation_name")
presentation = Presentation()

TITLE_AND_CONTENT = 1
BLANK = 6
###############################################################################
print("Set title", end="\n")

layout = presentation.slide_layouts[BLANK]
slide = presentation.slides.add_slide(layout)
###############################################################################
shapes = slide.shapes
left = Inches(4.0)
top = Inches(4.0)
width = Inches(3.0)
height = Inches(3.0)

shape = shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    left,
    top,
    width,
    height
)
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)
a = [(1, 1), [3, 5], [(3, 5), 4]]
###############################################################################
line = shape.line
line.color.rgb = RGBColor(255, 0, 0)
line.color.brightness = 0.5  # 50% lighter
# line.width = Pt(2.5)
# for number in range(10):
#
#     layout = presentation.slide_layouts[number]
#     slide = presentation.slides.add_slide(layout)
###############################################################################
save_to = parser.get("directory", "save_to")
presentation.save(f"{save_to}/{presentation_name}.pptx")

###############################################################################
finish_time = time.time()
print(f"Run time: {finish_time - start_time: .2f}", end="\n")
